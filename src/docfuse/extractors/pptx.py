"""Extracteur PPTX : .pptx.

CdC §8.3 — Texte des shapes, tableaux, notes d'orateur.
D-107 — Texte des graphiques, des SmartArt, du masque et des dispositions ;
marqueur de diapo honnête sur ce qui n'a pas pu être lu.
Détecte les images via ppt/media/* dans le ZIP.
"""

from __future__ import annotations

import logging
import re
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from docfuse.core.embedded_images import ImageBatch, build_image_tag
from docfuse.core.ocr.registry import resolve_ocr_engine
from docfuse.core.registry import register
from docfuse.extractors.base import Extractor, container_guard, error_result, file_type_for
from docfuse.models.extraction_result import ExtractedFile
from docfuse.models.file_status import FileStatus

logger = logging.getLogger(__name__)

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# `uri` de <a:graphicData> : ce qui distingue les sortes de <p:graphicFrame>.
_URI_CHART = _NS_C
_URI_DIAGRAM = _NS_DGM
_URI_TABLE = "http://schemas.openxmlformats.org/drawingml/2006/table"
_URI_OLE = "http://schemas.openxmlformats.org/presentationml/2006/ole"

# Éléments de <p:spTree> que python-pptx sait rendre sous forme de forme.
# Tout le reste (typiquement <mc:AlternateContent>, qui enveloppe l'encre
# manuscrite et les nouveautés PowerPoint) est ignoré **en silence** par la
# bibliothèque : on le repère nous-mêmes plutôt que de le perdre (D-107).
_TAGS_SPTREE_CONNUS = frozenset(
    f"{{{_NS_P}}}{nom}"
    for nom in ("nvGrpSpPr", "grpSpPr", "sp", "pic", "graphicFrame", "grpSp", "cxnSp")
)

_RE_DIAGRAM_DATA = re.compile(r"ppt/diagrams/data\d*\.xml")

# Au-delà, la liste des diapos concernées par une disposition est tronquée.
_MAX_DIAPOS_CITEES = 10


@register(".pptx")
class PptxExtractor(Extractor):
    """Extracteur PPTX via python-pptx + descente XML + détection media."""

    @classmethod
    def accepts(cls, path: Path) -> bool:
        return path.suffix.lower() == ".pptx"

    @classmethod
    def extract(cls, path: Path, relative_path: str, extract_images: bool = False) -> ExtractedFile:
        try:
            # D-089/D-093 : fichier protégé par mot de passe (conteneur OLE2)
            # ou « bombe zip » — garde partagée entre conteneurs (D-099).
            guard = container_guard(path, relative_path)
            if guard is not None:
                return guard

            from pptx import Presentation

            image_count = _count_media_images(path)

            # D-091 : OCR/export des images intégrées, seulement si utile
            # (export demandé ou OCR disponible) — sinon zéro coût ajouté.
            # D-098 : images collectées pendant le parcours (jetons), OCR de
            # tout le fichier en parallèle une fois, substitution par diapo.
            batch = ImageBatch(resolve_ocr_engine(), extract_images)

            prs = Presentation(str(path))
            slides: list[list[str]] = []
            non_lus_par_diapo: list[tuple[list[str], bool]] = []
            # D-107 : parties ppt/diagrams/dataN.xml effectivement rattachées
            # à une diapo, pour ne balayer en fin de course que les orphelines
            # (SmartArt posé sur un masque, une disposition, une page de
            # notes) — récupérées sans jamais être écrites deux fois.
            diagrammes_lus: set[str] = set()

            for i, slide in enumerate(prs.slides, 1):
                slide_text, non_lus, images_non_lues = _texte_diapo(
                    slide, i, relative_path, batch, diagrammes_lus
                )
                slides.append(slide_text)
                non_lus_par_diapo.append((non_lus, images_non_lues))

            batch.run()
            parts: list[str] = []

            # D-107 : le gabarit d'abord, une seule fois. Un bandeau de
            # classification vit sur le masque précisément pour valoir pour
            # toutes les diapos : le recopier sur chacune gonflerait le corpus
            # et fausserait le comptage de jetons.
            gabarit = _section_gabarit(prs)
            if gabarit:
                parts.append(gabarit)

            for i, slide_text in enumerate(slides, 1):
                resolved = batch.apply(slide_text)
                non_lus, images_non_lues = non_lus_par_diapo[i - 1]
                marqueur = _marqueur_diapo(i, not resolved, non_lus, images_non_lues)
                if marqueur:
                    resolved.append(marqueur)
                parts.append(f"## Diapo {i}\n\n" + "\n\n".join(resolved))

            orphelins = _section_diagrammes_orphelins(path, diagrammes_lus)
            if orphelins:
                parts.append(orphelins)

            slide_count = len(slides)
            embedded_images = batch.images
            text = "\n\n---\n\n".join(parts)

            return ExtractedFile(
                path=path,
                relative_path=relative_path,
                extension=file_type_for(path),
                file_type=file_type_for(path),
                size_bytes=path.stat().st_size,
                text=text,
                status=FileStatus.READY,
                image_count=image_count,
                page_count=slide_count,
                embedded_images=embedded_images,
            )
        except Exception as exc:
            logger.exception("Erreur extraction PPTX %s", path)
            return error_result(path, relative_path, exc)


def _texte_diapo(
    slide: Any,
    numero: int,
    relative_path: str,
    batch: ImageBatch,
    diagrammes_lus: set[str],
) -> tuple[list[str], list[str], bool]:
    """Texte d'une diapo et sources qu'on n'a pas su lire.

    Ce que la fonction rend permet au marqueur de fin de diapo de ne jamais
    affirmer faussement qu'il n'y a rien à lire (D-107).

    Returns:
        `(textes, sources non lues, au moins une image non analysée)`. Les
        images sont comptées à part des autres sources : elles sont déjà
        signalées au niveau du fichier (`image_count` → `FileStatus.IMAGES`
        dans le rapport), et un deck en compte des dizaines — les mentionner
        sur chaque diapo coûterait ~8 % de corpus pour une information déjà
        connue de l'auditeur.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    textes: list[str] = []
    non_lus: list[str] = []
    images_non_lues = False
    index_image = 0

    # D-074 : descend dans les formes groupées (GroupShape) — sans ça, tout
    # texte/tableau dans un groupe (schémas, diagrammes annotés — fréquents
    # dans les decks "corporate") est invisible : shape.has_text_frame /
    # shape.has_table renvoient False pour le conteneur groupe lui-même.
    for shape in _iter_shapes(slide.shapes):
        type_forme = _shape_type(shape)

        if type_forme == MSO_SHAPE_TYPE.PICTURE:
            if batch.active:
                index_image += 1
                token = _picture_token(shape, relative_path, numero, index_image, batch)
                if token:
                    textes.append(token)
            else:
                # Ni OCR ni export : l'image sort du corpus sans laisser de
                # trace. C'est exactement le cas où « aucun texte extractible »
                # mentait le plus souvent (D-107).
                images_non_lues = True
            continue

        if type_forme == MSO_SHAPE_TYPE.MEDIA:
            _ajouter(non_lus, "média audio/vidéo")
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = _clean_text(para.text)
                if text:
                    textes.append(text)

        if shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = [_clean_text(cell.text) for cell in row.cells]
                textes.append(" | ".join(cells))
            continue

        uri = _uri_graphic_data(shape)
        if uri is None or uri == _URI_TABLE:
            continue

        # D-107 : un <p:graphicFrame> n'a ni has_text_frame ni has_table.
        # Graphiques et SmartArt tombaient donc dans le vide, et la diapo
        # était déclarée « sans texte extractible » alors qu'elle porte des
        # catégories, des noms de séries, un organigramme nominatif.
        if uri == _URI_CHART:
            recupere = _texte_graphique(shape, numero)
            if recupere is None:
                _ajouter(non_lus, "graphique illisible")
            else:
                textes.extend(recupere)
        elif uri == _URI_DIAGRAM:
            recupere = _texte_smartart(shape, slide, numero, diagrammes_lus)
            if recupere is None:
                _ajouter(non_lus, "diagramme SmartArt illisible")
            else:
                textes.extend(recupere)
        elif uri == _URI_OLE:
            _ajouter(non_lus, "objet incorporé (OLE)")
        else:
            _ajouter(non_lus, f"objet graphique de type « {_fin_uri(uri)} »")

    recuperes, ignores = _elements_ignores_par_pptx(slide.shapes._spTree)
    textes.extend(recuperes)
    for etiquette in ignores:
        _ajouter(non_lus, etiquette)

    # Notes d'orateur (CdC §8.3)
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame
        if notes and _clean_text(notes.text):
            textes.append(f"[Notes] {_clean_text(notes.text)}")

    return textes, non_lus, images_non_lues


def _marqueur_diapo(numero: int, vide: bool, non_lus: list[str], images_non_lues: bool) -> str:
    """Marqueur de fin de diapo — jamais une affirmation fausse (D-107).

    L'ancien `[[DIAPO N: aucun texte extractible]]` disait à l'auditeur
    « diapo image, rien à lire » alors qu'un graphique, un organigramme
    SmartArt ou une image sans OCR venaient d'être jetés en silence. Quatre
    cas désormais distingués :

    - rien à signaler et du texte trouvé : pas de marqueur (inchangé) ;
    - diapo réellement dépourvue de forme porteuse de texte : on dit
      exactement ce qui a été inspecté, pas ce que la diapo contient ;
    - diapo sans texte lu mais porteuse de sources non lues : elles sont
      toutes nommées, images comprises — c'est le cas où l'ancien marqueur
      mentait, il doit être complet ;
    - diapo qui a produit du texte **et** porte une source non lue : la
      source est nommée, sauf s'il ne s'agit que d'images.

    Cette dernière exception est un arbitrage de volume, pas un oubli. Une
    image non analysée est déjà signalée au niveau du fichier (`image_count`
    → `FileStatus.IMAGES` dans le rapport et la liste). Mesuré sur huit
    decks réels sans Tesseract : la mention par diapo ajoutait 2 000
    caractères par fichier, soit ~8 % de corpus, sur des diapos qui portaient
    déjà leur texte — exactement le gonflement que l'on refuse au texte du
    masque. Une source *non-image* non lue reste signalée dans tous les cas :
    elle est rare (zéro occurrence sur ces mêmes huit decks) et peut cacher
    un bloc entier de contenu.
    """
    if vide and images_non_lues:
        non_lus = ["image non analysée (aucun moteur OCR disponible)", *non_lus]
    if non_lus:
        liste = " ; ".join(non_lus)
        if vide:
            return f"[[DIAPO {numero}: aucun texte lu — non analysé par DocFuse : {liste}]]"
        return f"[[DIAPO {numero}: également présent, non analysé par DocFuse : {liste}]]"
    if vide:
        return f"[[DIAPO {numero}: aucune forme porteuse de texte sur cette diapositive]]"
    return ""


def _section_gabarit(prs: Any) -> str:
    """Texte propre au masque et aux dispositions, rassemblé **une fois**.

    D-107 : `slide.shapes` ne voit pas ce qui vient de `slide_layout` ni de
    `slide_master` — un bandeau « DIFFUSION RESTREINTE » posé sur le masque
    (le cas d'usage même du masque) n'atteignait jamais le corpus.

    Pourquoi une section unique en tête du fichier et pas une répétition sur
    chaque diapo : ce texte est par construction commun à beaucoup de
    diapositives. Le recopier gonflerait le corpus proportionnellement au
    nombre de diapos et fausserait le compteur de contexte, pour une
    information qui ne varie pas. Les diapos concernées sont citées.

    Sont écartées les *marques de gabarit* : le texte d'invite des
    espaces réservés de titre/corps (« Cliquez pour modifier le style du
    titre »), la date et le numéro de diapo — ce sont des champs (`a:fld`),
    pas du contenu. L'espace réservé de pied de page est conservé : c'est là
    que vivent la moitié des mentions de classification.
    """
    masques: dict[str, Any] = {}
    dispositions: dict[str, tuple[Any, list[int]]] = {}
    nb_diapos = 0

    for i, slide in enumerate(prs.slides, 1):
        nb_diapos = i
        try:
            layout = slide.slide_layout
            cle = str(layout.part.partname)
            if cle not in dispositions:
                dispositions[cle] = (layout, [])
            dispositions[cle][1].append(i)
            master = layout.slide_master
            masques.setdefault(str(master.part.partname), master)
        except Exception:
            logger.warning("Gabarit de la diapo %d illisible", i, exc_info=True)

    lignes: list[str] = []
    vus: set[str] = set()

    for masque in masques.values():
        nom = _nom_gabarit(masque)
        etiquette = f"[Masque « {nom} »]" if nom else "[Masque]"
        for texte in _texte_gabarit(masque):
            if texte not in vus:
                vus.add(texte)
                lignes.append(f"{etiquette} {texte}")

    for layout, numeros in dispositions.values():
        nom = _nom_gabarit(layout)
        portee = _portee_diapos(numeros, nb_diapos)
        etiquette = f"[Disposition « {nom} » — {portee}]" if nom else f"[Disposition — {portee}]"
        for texte in _texte_gabarit(layout):
            if texte not in vus:
                vus.add(texte)
                lignes.append(f"{etiquette} {texte}")

    if not lignes:
        return ""
    entete = (
        "[[GABARIT: texte porté par le masque et les dispositions, commun à "
        "plusieurs diapositives — extrait ici une seule fois plutôt que "
        "répété sur chaque diapo]]"
    )
    return "## Gabarit\n\n" + entete + "\n\n" + "\n\n".join(lignes)


def _texte_gabarit(gabarit: Any) -> list[str]:
    """Texte réel d'un masque ou d'une disposition, marques de gabarit exclues."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    textes: list[str] = []
    for shape in _iter_shapes(gabarit.shapes):
        try:
            if shape.is_placeholder and shape.placeholder_format.type != PP_PLACEHOLDER.FOOTER:
                continue
        except Exception:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texte = _texte_paragraphe_sans_champs(para)
                if texte:
                    textes.append(texte)
        if shape.has_table:
            for row in shape.table.rows:
                ligne = " | ".join(_clean_text(cell.text) for cell in row.cells)
                if ligne.strip(" |"):
                    textes.append(ligne)
    return textes


def _texte_paragraphe_sans_champs(para: Any) -> str:
    """Texte d'un paragraphe sans le contenu des champs `a:fld`.

    `paragraph.text` de python-pptx inclut la valeur figée des champs — sur
    un masque, c'est la date de création du modèle (« 1/27/13 ») et le
    littéral « ‹#› » du numéro de diapo. Du bruit pur dans le corpus.
    `paragraph.runs` les exclut mais perd aussi les sauts de ligne manuels
    (`a:br`), ce qui collerait deux mots : on reconstruit donc à la main.
    """
    morceaux: list[str] = []
    for enfant in para._p:
        if enfant.tag == f"{{{_NS_A}}}r":
            morceaux.append("".join(t.text or "" for t in enfant.iter(f"{{{_NS_A}}}t")))
        elif enfant.tag == f"{{{_NS_A}}}br":
            morceaux.append("\n")
    return _clean_text("".join(morceaux))


def _nom_gabarit(gabarit: Any) -> str:
    try:
        return str(gabarit.name or "").strip()
    except Exception:
        return ""


def _portee_diapos(numeros: list[int], nb_diapos: int) -> str:
    if nb_diapos and len(numeros) == nb_diapos:
        return "toutes les diapositives"
    if len(numeros) > _MAX_DIAPOS_CITEES:
        debut = ", ".join(str(n) for n in numeros[:_MAX_DIAPOS_CITEES])
        return f"diapos {debut}… ({len(numeros)} diapositives)"
    return "diapos " + ", ".join(str(n) for n in numeros)


def _texte_graphique(shape: Any, numero: int) -> list[str] | None:
    """Titre, noms de séries, catégories et étiquettes d'un graphique.

    Lu dans le XML de la partie graphique plutôt que par l'API `chart` de
    python-pptx : celle-ci ne couvre pas tous les types de graphique
    (`plots`/`categories` lèvent sur certains), alors que les caches de
    chaînes ont la même forme pour tous. Les valeurs numériques
    (`c:numCache`) sont volontairement laissées de côté : ce sont des
    nombres sans étiquette, du bruit pour un classement de sensibilité.

    Returns:
        Les textes trouvés (liste éventuellement vide si le graphique n'a
        aucun libellé), ou `None` si la partie est illisible.
    """
    try:
        blob = shape.chart.part.blob
        racine = ElementTree.fromstring(blob)
    except Exception:
        logger.warning("Graphique illisible sur la diapo %d", numero, exc_info=True)
        return None

    bruts = [elt.text or "" for elt in racine.iter(f"{{{_NS_A}}}t")]
    for chemin in (
        f".//{{{_NS_C}}}strCache/{{{_NS_C}}}pt/{{{_NS_C}}}v",
        f".//{{{_NS_C}}}strLit/{{{_NS_C}}}pt/{{{_NS_C}}}v",
        f".//{{{_NS_C}}}multiLvlStrRef//{{{_NS_C}}}pt/{{{_NS_C}}}v",
        f".//{{{_NS_C}}}tx/{{{_NS_C}}}v",
    ):
        bruts.extend(elt.text or "" for elt in racine.findall(chemin))
    return _dedupe(_clean_text(t) for t in bruts)


def _texte_smartart(
    shape: Any, slide: Any, numero: int, diagrammes_lus: set[str]
) -> list[str] | None:
    """Texte d'un SmartArt, lu dans `ppt/diagrams/dataN.xml`.

    python-pptx n'expose rien d'un SmartArt : la forme est un
    `<p:graphicFrame>` dont le `<dgm:relIds r:dm="...">` pointe une partie
    « diagramData » du paquet. Un organigramme y est du nominatif pur — la
    perte la plus coûteuse de l'ancien code.

    Returns:
        Les textes du modèle de données, ou `None` si la partie liée est
        introuvable ou illisible.
    """
    try:
        rel_ids = shape._element.find(f".//{{{_NS_DGM}}}relIds")
        rid = None if rel_ids is None else rel_ids.get(f"{{{_NS_R}}}dm")
        if not rid:
            return None
        part = slide.part.rels[rid].target_part
        diagrammes_lus.add(str(part.partname))
        return _texte_data_model(part.blob)
    except Exception:
        logger.warning("SmartArt illisible sur la diapo %d", numero, exc_info=True)
        return None


def _texte_data_model(blob: bytes) -> list[str]:
    """Paragraphes d'un modèle de données SmartArt (`dgm:dataModel`)."""
    racine = ElementTree.fromstring(blob)
    bruts: list[str] = []
    for noeud in racine.iter(f"{{{_NS_DGM}}}t"):
        for para in noeud.iter(f"{{{_NS_A}}}p"):
            bruts.append("".join(elt.text or "" for elt in para.iter(f"{{{_NS_A}}}t")))
    return _dedupe(_clean_text(t) for t in bruts)


def _section_diagrammes_orphelins(path: Path, diagrammes_lus: set[str]) -> str:
    """SmartArt du paquet qu'aucune diapo n'a réclamé (masque, disposition,
    page de notes) — récupérés en fin de corpus plutôt que perdus (D-107).
    Les parties déjà lues par une diapo ne sont pas répétées."""
    lignes: list[str] = []
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            for nom in sorted(n for n in zf.namelist() if _RE_DIAGRAM_DATA.fullmatch(n)):
                if f"/{nom}" in diagrammes_lus:
                    continue
                try:
                    lignes.extend(_texte_data_model(zf.read(nom)))
                except Exception:
                    logger.warning("Diagramme %s illisible", nom, exc_info=True)
    except Exception:
        logger.warning("Balayage des diagrammes de %s impossible", path, exc_info=True)

    if not lignes:
        return ""
    entete = (
        "[[DIAGRAMMES: SmartArt du fichier qui n'ont pas pu être rattachés à "
        "une diapositive précise (posés sur un masque, une disposition ou une "
        "page de notes) — restitués ici pour ne pas les perdre]]"
    )
    return "## Diagrammes hors diapositive\n\n" + entete + "\n\n" + "\n\n".join(_dedupe(lignes))


def _elements_ignores_par_pptx(sp_tree: Any) -> tuple[list[str], list[str]]:
    """Contenu de `<p:spTree>` que python-pptx n'énumère pas comme forme.

    La bibliothèque saute sans un mot tout élément qu'elle ne connaît pas —
    au premier rang `<mc:AlternateContent>`, l'enveloppe que PowerPoint pose
    autour de l'encre manuscrite et de ses nouveautés. On récupère le texte
    de la branche retenue (`mc:Choice` de préférence, `mc:Fallback` sinon —
    jamais les deux, elles portent le même contenu) et, à défaut de texte,
    on signale l'élément plutôt que de le passer sous silence.

    Returns:
        (textes récupérés, étiquettes des sources restées non lues).
    """
    textes: list[str] = []
    non_lus: list[str] = []
    a_visiter = [sp_tree]
    while a_visiter:
        for element in a_visiter.pop():
            if element.tag == f"{{{_NS_P}}}grpSp":
                a_visiter.append(element)
                continue
            if element.tag in _TAGS_SPTREE_CONNUS:
                continue
            branche = _branche_retenue(element)
            trouves = _dedupe(
                _clean_text("".join(t.text or "" for t in para.iter(f"{{{_NS_A}}}t")))
                for para in branche.iter(f"{{{_NS_A}}}p")
            )
            if trouves:
                textes.extend(trouves)
            else:
                _ajouter(non_lus, f"élément « {_nom_local(element.tag)} » non lu")
    return textes, non_lus


def _branche_retenue(element: Any) -> Any:
    """Pour un `<mc:AlternateContent>`, la branche à lire ; l'élément lui-même
    sinon. Sans ce choix, `mc:Choice` et `mc:Fallback` livreraient deux fois
    le même texte."""
    for suffixe in ("Choice", "Fallback"):
        for enfant in element:
            if _nom_local(enfant.tag) == suffixe:
                return enfant
    return element


def _nom_local(tag: Any) -> str:
    return str(tag).rsplit("}", 1)[-1]


def _fin_uri(uri: str) -> str:
    return uri.rstrip("/").rsplit("/", 1)[-1] or uri


def _uri_graphic_data(shape: Any) -> str | None:
    """`uri` de `<a:graphicData>` — la seule chose qui dise ce qu'un
    `<p:graphicFrame>` contient. `None` si la forme n'en est pas un."""
    try:
        return str(shape._element.graphic.graphicData.uri)
    except Exception:
        return None


def _shape_type(shape: Any) -> Any:
    """`shape.shape_type` sans jamais lever : la propriété n'est pas définie
    pour tous les `<p:graphicFrame>` (un SmartArt renvoie `None`) et une
    exception ici ferait basculer tout le fichier en ERREUR."""
    try:
        return shape.shape_type
    except Exception:
        return None


def _ajouter(liste: list[str], valeur: str) -> None:
    """Ajoute une étiquette de source non lue si elle n'y est pas déjà — une
    diapo à quinze photos ne doit produire qu'une mention."""
    if valeur not in liste:
        liste.append(valeur)


def _dedupe(valeurs: Any) -> list[str]:
    """Textes non vides, sans doublon, dans l'ordre d'apparition."""
    vus: set[str] = set()
    resultat: list[str] = []
    for valeur in valeurs:
        if valeur and valeur not in vus:
            vus.add(valeur)
            resultat.append(valeur)
    return resultat


def _clean_text(text: str) -> str:
    """Normalise le texte d'une forme PPTX (D-096).

    python-pptx rend un saut de ligne manuel (`<a:br/>`, Maj+Entrée dans
    PowerPoint) comme `\\x0b` (tabulation verticale) dans `.text`. Sans
    conversion, ce caractère de contrôle finit tel quel dans le Markdown
    (compté par les tokenizers, rendu en glyphe inconnu dans le PDF).
    """
    return text.replace("\x0b", "\n").strip()


def _picture_token(
    shape: Any, relative_path: str, slide_no: int, index: int, batch: ImageBatch
) -> str:
    """Enregistre l'image d'une forme dans le lot et renvoie son jeton de
    position (D-091, D-098). N'échoue jamais : une image illisible est
    ignorée plutôt que de faire échouer toute l'extraction de la diapo."""
    try:
        image = shape.image
        data = image.blob
        ext = image.ext
    except Exception:
        logger.warning("Lecture de l'image slide %d échouée", slide_no, exc_info=True)
        return ""
    return batch.add(build_image_tag(relative_path, f"slide{slide_no}", index, ext), data)


def _iter_shapes(shapes: Any) -> Any:
    """Parcourt les formes d'une diapo récursivement, en descendant dans les
    formes groupées (GroupShape, D-074) — la notion de groupe est elle-même
    récursive (un groupe peut contenir un groupe)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if _shape_type(shape) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _count_media_images(path: Path) -> int:
    """Compte les images dans ppt/media/ du ZIP PPTX."""
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            return sum(1 for n in zf.namelist() if n.startswith("ppt/media/"))
    except Exception:
        return 0
