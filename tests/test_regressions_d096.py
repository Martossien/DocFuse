"""Non-régression de l'audit D-096 (lot 1 : contenu perdu / plantage total).

Chaque test reproduit exactement le cas constaté pendant l'audit (avec la
vraie bibliothèque du format, jamais un mock) puis vérifie le correctif.
Voir docs/journal-decisions.md — D-096.
"""

from __future__ import annotations

import io
import zipfile
from email.message import EmailMessage
from pathlib import Path

import pytest

from docfuse.core.orchestrator import generate_corpus, run_analysis
from docfuse.models.file_status import FileStatus

_ODF_NS = (
    'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
    'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" '
    'xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"'
)


def _odf(path: Path, body: str, mime: str) -> Path:
    xml = (
        f'<?xml version="1.0"?><office:document-content {_ODF_NS}>'
        f"<office:body>{body}</office:body></office:document-content>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("mimetype", mime)
        zf.writestr("content.xml", xml)
    return path


class TestOrchestratorDuplicateRemoval:
    def test_removing_duplicate_original_promotes_copy(self, tmp_path: Path) -> None:
        """Retirer l'original d'un doublon laissait seulement la note
        « identique à … » dans le corpus : le contenu réel disparaissait."""
        text = "Contenu strictement identique repete pour depasser le seuil minimum requis."
        (tmp_path / "1_original.txt").write_text(text, encoding="utf-8")
        (tmp_path / "2_copie.txt").write_text(text, encoding="utf-8")
        (tmp_path / "3_copie.txt").write_text(text, encoding="utf-8")
        result = run_analysis(tmp_path, context_limit=128000)

        assert result.remove_file(tmp_path / "1_original.txt", reason="test")

        by_name = {f.relative_path: f for f in result.files}
        assert by_name["2_copie.txt"].text == text
        assert "duplicate_of" not in by_name["2_copie.txt"].extra_metadata
        assert by_name["3_copie.txt"].extra_metadata["duplicate_of"] == "2_copie.txt"
        assert result.total.tokens_estimated > 0


class TestInventoryRobustness:
    def test_mtime_sort_survives_broken_symlink(self, tmp_path: Path) -> None:
        (tmp_path / "ok.txt").write_text("ok", encoding="utf-8")
        (tmp_path / "broken.txt").symlink_to(tmp_path / "nonexistent-target.txt")

        result = run_analysis(tmp_path, context_limit=128000, sort="mtime")

        statuses = {f.relative_path: f.status for f in result.files}
        assert statuses["broken.txt"] is FileStatus.ERROR
        assert statuses["ok.txt"].is_extracted()

    def test_pruned_directories_are_listed_as_ignored(self, tmp_path: Path) -> None:
        """`build/`, `dist/`… étaient élagués sans apparaître dans le rapport."""
        (tmp_path / "ok.txt").write_text("ok", encoding="utf-8")
        (tmp_path / "build").mkdir()
        (tmp_path / "build" / "permis_de_construire.txt").write_text("x", encoding="utf-8")

        result = run_analysis(tmp_path, context_limit=128000)

        assert [f.relative_path for f in result.files] == ["ok.txt"]
        assert any(p.name == "build" for p, _ in result.ignored)


class TestXlsxChartsheet:
    def test_chartsheet_does_not_lose_data_sheets(self, tmp_path: Path) -> None:
        import openpyxl
        from openpyxl.chart import BarChart, Reference

        from docfuse.extractors.xlsx import XlsxExtractor

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        for i in range(1, 6):
            ws.cell(row=i, column=1, value=i * 10)
        chart_sheet = wb.create_chartsheet("Graphique")
        chart = BarChart()
        chart.add_data(Reference(ws, min_col=1, min_row=1, max_row=5))
        chart_sheet.add_chart(chart)
        f = tmp_path / "chart.xlsx"
        wb.save(str(f))

        result = XlsxExtractor.extract(f, "chart.xlsx")

        assert result.status is FileStatus.READY
        assert "50" in result.text
        assert "Feuille : Graphique" in result.text
        assert "graphique" in result.text.lower()
        assert result.page_count == 2


class TestDocxStructure:
    def test_linked_header_emitted_once(self, tmp_path: Path) -> None:
        from docx import Document
        from docx.enum.section import WD_SECTION

        from docfuse.extractors.docx import DocxExtractor

        doc = Document()
        doc.sections[0].header.paragraphs[0].text = "EN-TETE UNIQUE"
        doc.add_paragraph("Corps.")
        for _ in range(3):
            doc.add_section(WD_SECTION.NEW_PAGE)
            doc.add_paragraph("Suite.")
        f = tmp_path / "sections.docx"
        doc.save(str(f))

        result = DocxExtractor.extract(f, "sections.docx")
        assert result.text.count("EN-TETE UNIQUE") == 1

    def test_textbox_emitted_once_without_gluing_runs(self, tmp_path: Path) -> None:
        from docx import Document

        from docfuse.extractors.docx import DocxExtractor

        f = tmp_path / "tb.docx"
        doc = Document()
        doc.add_paragraph("Corps normal.")
        doc.save(str(f))
        textbox = (
            "<w:p><w:r><w:pict><v:shape><v:textbox><w:txbxContent>"
            "<w:p><w:r><w:t>Hello</w:t></w:r><w:r><w:t xml:space='preserve'> world</w:t></w:r></w:p>"
            "<w:p><w:r><w:t>Second para</w:t></w:r></w:p>"
            "</w:txbxContent></v:textbox></v:shape></w:pict></w:r></w:p>"
        )
        with zipfile.ZipFile(f) as zin:
            items = zin.infolist()
            contents = {i.filename: zin.read(i.filename) for i in items}
        xml = (
            contents["word/document.xml"]
            .decode("utf-8")
            .replace("</w:body>", textbox + "</w:body>")
        )
        contents["word/document.xml"] = xml.encode("utf-8")
        with zipfile.ZipFile(f, "w") as zout:
            for item in items:
                zout.writestr(item, contents[item.filename])

        result = DocxExtractor.extract(f, "tb.docx")

        assert result.text.count("Hello") == 1
        assert "Hello world\nSecond para" in result.text


class TestPptxLineBreak:
    def test_manual_line_break_is_newline_not_vertical_tab(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        from docfuse.extractors.pptx import PptxExtractor

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "Ligne A"
        box.text_frame.paragraphs[0].add_line_break()
        box.text_frame.paragraphs[0].add_run().text = "Ligne B"
        f = tmp_path / "br.pptx"
        prs.save(str(f))

        result = PptxExtractor.extract(f, "br.pptx")
        assert "\x0b" not in result.text
        assert "Ligne A\nLigne B" in result.text


class TestPdfGarbagePages:
    def test_blank_if_garbage_only_for_ocr_kind(self) -> None:
        from docfuse.extractors.pdf import PageKind, _blank_if_garbage

        assert _blank_if_garbage(PageKind.OCR, "(cid:12)(cid:13)") == ""
        assert _blank_if_garbage(PageKind.OCR, "texte court réel") == "texte court réel"
        assert _blank_if_garbage(PageKind.MIXED, "(cid:12)") == "(cid:12)"


class TestWriters:
    def test_pdf_writer_escapes_source_header(self, tmp_path: Path) -> None:
        (tmp_path / "a<b>.txt").write_text("Ligne 1\nLigne 2\n", encoding="utf-8")
        result = run_analysis(tmp_path, context_limit=128000)

        assert generate_corpus(result, tmp_path / "corpus.pdf")
        assert (tmp_path / "corpus.pdf").stat().st_size > 0

    def test_markdown_crlf_mode_has_no_bare_lf(self, tmp_path: Path) -> None:
        from docfuse.output.markdown_writer import write_markdown_corpus

        (tmp_path / "a.txt").write_text("l1\nl2\r\nl3\rl4\n", encoding="utf-8", newline="")
        result = run_analysis(tmp_path, context_limit=128000)
        out = tmp_path / "corpus.md"
        write_markdown_corpus(result, out, 0.15, line_ending="crlf")

        data = out.read_bytes()
        assert data.count(b"\r\n") > 0
        assert data.count(b"\n") == data.count(b"\r\n")
        assert data.replace(b"\r\n", b"").count(b"\r") == 0


class TestConfigRobustness:
    def test_non_numeric_context_limit_falls_back_to_defaults(self, tmp_path: Path) -> None:
        from docfuse.config import load_config

        cfg = tmp_path / "c.json"
        cfg.write_text('{"context_limit": "abc"}', encoding="utf-8")
        assert load_config(cfg).context_limit == 128_000

    def test_exclude_globs_string_is_one_pattern(self, tmp_path: Path) -> None:
        from docfuse.config import load_config

        cfg = tmp_path / "c.json"
        cfg.write_text('{"exclude_globs": "*.log"}', encoding="utf-8")
        assert load_config(cfg).exclude_globs == ["*.log"]

    def test_boolean_strings_are_parsed_strictly(self, tmp_path: Path) -> None:
        from docfuse.config import load_config

        cfg = tmp_path / "c.json"
        cfg.write_text('{"recursive": "false"}', encoding="utf-8")
        assert load_config(cfg).recursive is False

    def test_validate_is_enforced_by_cli(self, tmp_path: Path) -> None:
        from docfuse.cli import main

        cfg = tmp_path / "c.json"
        cfg.write_text('{"context_limit": -5}', encoding="utf-8")
        (tmp_path / "a.txt").write_text("x", encoding="utf-8")
        code = main(["-i", str(tmp_path), "--config", str(cfg), "-o", str(tmp_path / "c.md")])
        assert code == 1


class TestHtmlStructure:
    def test_div_container_keeps_headings_tables_lists(self, tmp_path: Path) -> None:
        from docfuse.extractors.html import HtmlExtractor

        f = tmp_path / "div.html"
        f.write_text(
            "<html><body><div><h1>Titre</h1><p>Para un.</p>"
            "<table><tr><td>A</td><td>B</td></tr></table>"
            "<ul><li>x</li><li>y</li></ul></div></body></html>",
            encoding="utf-8",
        )
        result = HtmlExtractor.extract(f, "div.html")
        assert "# Titre" in result.text
        assert "| A | B |" in result.text
        assert "- x" in result.text

    def test_inline_formatting_does_not_glue_words(self, tmp_path: Path) -> None:
        from docfuse.extractors.html import HtmlExtractor

        f = tmp_path / "glue.html"
        f.write_text(
            "<html><body><h1>Hello <b>World</b> again</h1>"
            "<ul><li>Item <a>link</a> more</li></ul>"
            "<table><tr><td>Total <b>100</b> EUR</td></tr></table>"
            "<pre>def f():\n    return 1</pre></body></html>",
            encoding="utf-8",
        )
        result = HtmlExtractor.extract(f, "glue.html")
        assert "# Hello World again" in result.text
        assert "- Item link more" in result.text
        assert "| Total 100 EUR |" in result.text
        assert "def f():\n    return 1" in result.text


class TestOdfStructure:
    def test_spans_and_whitespace_elements(self, tmp_path: Path) -> None:
        from docfuse.extractors.odf import OdfExtractor

        f = _odf(
            tmp_path / "span.odt",
            "<office:text><text:p>Bonjour <text:span>monde</text:span> entier</text:p>"
            '<text:p>col1<text:tab/>col2<text:s text:c="3"/>col3</text:p></office:text>',
            "application/vnd.oasis.opendocument.text",
        )
        result = OdfExtractor.extract(f, "span.odt")
        assert "Bonjour monde entier" in result.text
        assert "col1\tcol2 col3" in result.text

    def test_section_content_is_not_dropped(self, tmp_path: Path) -> None:
        from docfuse.extractors.odf import OdfExtractor

        f = _odf(
            tmp_path / "section.odt",
            "<office:text><text:p>Avant</text:p><text:section><text:h>TITRE_SECTION</text:h>"
            "<text:p>PARA_SECTION</text:p></text:section><text:p>Apres</text:p></office:text>",
            "application/vnd.oasis.opendocument.text",
        )
        result = OdfExtractor.extract(f, "section.odt")
        assert result.text == "Avant\nTITRE_SECTION\nPARA_SECTION\nApres"

    def test_ods_keeps_rows_and_columns(self, tmp_path: Path) -> None:
        from docfuse.extractors.odf import OdfExtractor

        f = _odf(
            tmp_path / "t.ods",
            "<office:spreadsheet><table:table table:name='F1'>"
            "<table:table-row><table:table-cell><text:p>Nom</text:p></table:table-cell>"
            "<table:table-cell><text:p>Age</text:p></table:table-cell></table:table-row>"
            "<table:table-row><table:table-cell><text:p>Bob</text:p></table:table-cell>"
            "<table:table-cell><text:p>42</text:p></table:table-cell></table:table-row>"
            "</table:table></office:spreadsheet>",
            "application/vnd.oasis.opendocument.spreadsheet",
        )
        result = OdfExtractor.extract(f, "t.ods")
        assert "### Feuille : F1" in result.text
        assert "Nom | Age" in result.text
        assert "Bob | 42" in result.text


class TestEmlRobustness:
    def test_unknown_charset_does_not_fail_whole_email(self, tmp_path: Path) -> None:
        from docfuse.extractors.eml import EmlExtractor

        f = tmp_path / "unk.eml"
        f.write_bytes(
            b"From: a@b.c\r\nTo: d@e.f\r\nSubject: T\r\n"
            b"Content-Type: text/plain; charset=unknown-8bit\r\n\r\nCorps avec accent \xe9\r\n"
        )
        result = EmlExtractor.extract(f, "unk.eml")
        assert result.status is FileStatus.READY
        assert "Corps avec accent" in result.text

    def test_attachment_never_replaces_body_and_is_listed(self, tmp_path: Path) -> None:
        from docfuse.extractors.eml import EmlExtractor

        msg = EmailMessage()
        msg["From"] = "a@b.c"
        msg["To"] = "d@e.f"
        msg["Cc"] = "cc@x.y"
        msg["Subject"] = "S"
        msg.set_content("<p>CORPS_HTML_REEL</p>", subtype="html")
        msg.add_attachment(
            b"CONTENU_PIECE_JOINTE", maintype="text", subtype="plain", filename="notes.txt"
        )
        f = tmp_path / "att.eml"
        f.write_bytes(bytes(msg))

        result = EmlExtractor.extract(f, "att.eml")
        assert "CORPS_HTML_REEL" in result.text
        assert "CONTENU_PIECE_JOINTE" not in result.text
        assert "notes.txt" in result.text
        assert "cc@x.y" in result.text


class TestEpubSpine:
    def test_percent_encoded_href_and_missing_item_are_handled(self, tmp_path: Path) -> None:
        from docfuse.extractors.epub import EpubExtractor

        f = tmp_path / "pct.epub"
        with zipfile.ZipFile(f, "w") as zf:
            zf.writestr("mimetype", "application/epub+zip")
            zf.writestr(
                "META-INF/container.xml",
                '<?xml version="1.0"?><container xmlns="urn:oasis:names:tc:opendocument:xmlns:container" version="1.0">'
                '<rootfiles><rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/></rootfiles></container>',
            )
            zf.writestr(
                "OEBPS/content.opf",
                '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="2.0">'
                '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/"/>'
                '<manifest><item id="c1" href="chap%201.xhtml" media-type="application/xhtml+xml"/>'
                '<item id="c2" href="absent.xhtml" media-type="application/xhtml+xml"/></manifest>'
                '<spine><itemref idref="c1"/><itemref idref="c2"/></spine></package>',
            )
            zf.writestr("OEBPS/chap 1.xhtml", "<html><body><p>CHAPITRE_UN</p></body></html>")

        result = EpubExtractor.extract(f, "pct.epub")
        assert "CHAPITRE_UN" in result.text
        assert "absent.xhtml" in result.extra_metadata["epub_skipped_items"]


class TestXmlDeclaredEncoding:
    def test_declaration_wins_and_comments_survive(self, tmp_path: Path) -> None:
        from docfuse.extractors.xml_json import XmlExtractor

        f = tmp_path / "cp1251.xml"
        f.write_bytes(
            '<?xml version="1.0" encoding="windows-1251"?><t><!-- COMMENTAIRE --><a>Привет мир</a></t>'.encode(
                "cp1251"
            )
        )
        result = XmlExtractor.extract(f, "cp1251.xml")
        assert result.status is FileStatus.READY
        assert result.encoding == "windows-1251"
        assert "Привет мир" in result.text
        assert "COMMENTAIRE" in result.text


class TestSmallWholeFileFailures:
    def test_rtf_undefined_cp1252_byte_degrades_locally(self, tmp_path: Path) -> None:
        from docfuse.extractors.rtf import RtfExtractor

        f = tmp_path / "bad.rtf"
        f.write_bytes(b"{\\rtf1\\ansi a\\'81b}")
        result = RtfExtractor.extract(f, "bad.rtf")
        assert result.status is FileStatus.READY
        assert result.text.startswith("a")
        assert result.text.endswith("b")

    def test_csv_field_over_default_limit(self, tmp_path: Path) -> None:
        from docfuse.extractors.csv_tsv import CsvTsvExtractor

        f = tmp_path / "big.csv"
        f.write_text("id,blob\n1," + "x" * 200_000 + "\n", encoding="utf-8")
        result = CsvTsvExtractor.extract(f, "big.csv")
        assert result.status is FileStatus.READY
        assert len(result.text) > 200_000


class TestReportAlignment:
    def test_misaligned_estimates_fail_loudly(self, tmp_path: Path) -> None:
        from docfuse.core.report import generate_markdown_report

        (tmp_path / "a.txt").write_text("x", encoding="utf-8")
        result = run_analysis(tmp_path, context_limit=128000)
        with pytest.raises(ValueError, match="alignés"):
            generate_markdown_report(
                result.files, result.ignored, 128000, 0.15, 1, 1, tmp_path / "r.md", estimates=[]
            )


class TestGuiDragAndDrop:
    def test_tkdnd_package_is_loaded_into_the_interpreter(self) -> None:
        """Le paquet Tcl doit être chargé (`TkinterDnD.require`) — sans ça
        `drop_target_register` échoue toujours (D-096)."""
        source = (Path(__file__).resolve().parent.parent / "src/docfuse/gui.py").read_text(
            encoding="utf-8"
        )
        assert "TkinterDnD.require(" in source

    def test_specs_bundle_tkdnd_tcl_library(self) -> None:
        root = Path(__file__).resolve().parent.parent
        for spec in ("CorpusOne.spec", "CorpusOne-OCR.spec"):
            assert 'collect_data_files("tkinterdnd2")' in (root / spec).read_text(encoding="utf-8")


def _png_bytes() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (10, 10), "white").save(buf, format="PNG")
    return buf.getvalue()
