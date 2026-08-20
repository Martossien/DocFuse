#!/usr/bin/env python3
"""Script de recette DocFuse (CdC S21.4).

Lance une serie de tests fonctionnels sur le jeu de fichiers de test anonymise.
Verifie que les cas d'acceptation du CdC S19 fonctionnent correctement.

Usage:
    python tests/recette/run_recette.py
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

# ------------------------------------------------------------------------------
# Configuration
# ------------------------------------------------------------------------------

RECETTE_DIR = Path(__file__).resolve().parent
DOSSIER_MIXTE = RECETTE_DIR / "dossier_mixte"
DOSSIER_BLOCAGE = RECETTE_DIR / "dossier_blocage"
DOSSIER_IMAGES = RECETTE_DIR / "dossier_images"

# ------------------------------------------------------------------------------
# Generation des fichiers de test (CI n'a pas les fixtures versionnees)
# ------------------------------------------------------------------------------


def _generate_test_files() -> None:
    """Genere les fichiers de test s'ils n'existent pas (CI)."""
    import json
    from email.message import EmailMessage

    if DOSSIER_MIXTE.exists() and any(DOSSIER_MIXTE.iterdir()):
        return  # Deja genere

    DOSSIER_MIXTE.mkdir(parents=True, exist_ok=True)
    (DOSSIER_MIXTE / "note.txt").write_text(
        "Ceci est une note interne avec suffisamment de texte pour depasser le seuil de 80 caracteres.\n",
        encoding="utf-8",
    )
    (DOSSIER_MIXTE / "lisez-moi.md").write_text(
        "# README\n\nCe document explique le fonctionnement avec assez de caracteres.\n",
        encoding="utf-8",
    )
    (DOSSIER_MIXTE / "page.html").write_text(
        "<html><body><h1>Page de Test</h1><p>Contenu de la page avec assez de caracteres pour eviter l alerte.</p></body></html>",
        encoding="utf-8",
    )
    (DOSSIER_MIXTE / "config.json").write_text(
        json.dumps({"projet": "Test", "version": "1.0"}), encoding="utf-8"
    )
    (DOSSIER_MIXTE / "donnees.csv").write_text(
        "nom;valeur;description\nTest1;42;Premiere ligne\nTest2;99;Deuxieme ligne\n",
        encoding="utf-8",
    )
    (DOSSIER_MIXTE / "app.exe").write_bytes(b"\x4d\x5a\x90\x00")
    (DOSSIER_MIXTE / "~$locked.docx").write_bytes(b"\x00")
    (DOSSIER_MIXTE / "photo.jpg").write_bytes(b"\xff\xd8\xff\xe0\x00\x10JFIF")

    # DOCX
    try:
        from docx import Document

        doc = Document()
        doc.add_heading("Document Test", 0)
        doc.add_paragraph(
            "Paragraphe de test avec assez de caracteres pour eviter l alerte de pauvrete."
        )
        doc.save(str(DOSSIER_MIXTE / "rapport.docx"))
    except Exception:
        pass

    # PPTX
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Diapo Test"
        slide.placeholders[1].text = "Contenu avec assez de caracteres pour eviter l alerte."
        prs.save(str(DOSSIER_MIXTE / "slides.pptx"))
    except Exception:
        pass

    # XLSX
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Nom"
        ws["B1"] = "Valeur"
        ws["A2"] = "Test avec du texte"
        ws["B2"] = "42"
        wb.save(str(DOSSIER_MIXTE / "donnees.xlsx"))
    except Exception:
        pass

    # PDF
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate

        doc = SimpleDocTemplate(str(DOSSIER_MIXTE / "document.pdf"), pagesize=A4)
        styles = getSampleStyleSheet()
        doc.build(
            [
                Paragraph(
                    "Texte PDF de test avec assez de caracteres pour eviter l alerte.",
                    styles["Normal"],
                )
            ]
        )
    except Exception:
        pass

    # EML
    msg = EmailMessage()
    msg["Subject"] = "Test EML"
    msg["From"] = "test@example.com"
    msg["To"] = "recipient@example.com"
    msg.set_content("Corps de l email avec assez de caracteres pour eviter l alerte de pauvrete.")
    (DOSSIER_MIXTE / "email.eml").write_bytes(bytes(msg))

    # RTF
    (DOSSIER_MIXTE / "document.rtf").write_text(
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Arial;}} \f0\fs24 Texte RTF de test avec assez de caracteres pour eviter l alerte.}",
        encoding="latin-1",
    )

    # ODT
    try:
        import zipfile

        content = (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
            'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
            "<office:body><office:text>"
            "<text:p>Texte ODT de test avec assez de caracteres pour eviter l alerte.</text:p>"
            "</office:text></office:body></office:document-content>"
        )
        with zipfile.ZipFile(str(DOSSIER_MIXTE / "open-doc.odt"), "w") as zf:
            zf.writestr("mimetype", "application/vnd.oasis.opendocument.text")
            zf.writestr("content.xml", content)
    except Exception:
        pass

    # Sous-dossier
    (DOSSIER_MIXTE / "sous-dossier").mkdir(exist_ok=True)
    (DOSSIER_MIXTE / "sous-dossier" / "note-profonde.md").write_text(
        "# Note profonde\n\nTexte dans un sous-dossier avec assez de caracteres.\n",
        encoding="utf-8",
    )

    # Dossier blocage
    DOSSIER_BLOCAGE.mkdir(parents=True, exist_ok=True)
    (DOSSIER_BLOCAGE / "gros.txt").write_text("A" * 10000, encoding="utf-8")

    # Dossier images
    DOSSIER_IMAGES.mkdir(parents=True, exist_ok=True)
    (DOSSIER_IMAGES / "note.txt").write_text(
        "Texte avec assez de caracteres pour eviter l alerte de pauvrete dans DocFuse.\n",
        encoding="utf-8",
    )


# ------------------------------------------------------------------------------
# Utilitaires
# ------------------------------------------------------------------------------


class Result:
    """Resultat d'un test de recette."""

    def __init__(self, name: str) -> None:
        self.name = name
        self.passed = False
        self.message = ""

    def ok(self, msg: str = "") -> None:
        self.passed = True
        self.message = msg

    def fail(self, msg: str) -> None:
        self.passed = False
        self.message = msg


def run_cli(args: list[str]) -> tuple[int, str, str]:
    """Lance la CLI et retourne (code retour, stdout, stderr)."""
    cmd = [sys.executable, "-m", "docfuse"] + args
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
    return proc.returncode, proc.stdout, proc.stderr


# ------------------------------------------------------------------------------
# Tests de recette (CdC S19)
# ------------------------------------------------------------------------------


def test_dossier_mixte_md() -> Result:
    """S19.2 -- Dossier mixte -> un MD unique, chaque source identifiable."""
    r = Result("Dossier mixte -> Markdown")
    output = DOSSIER_MIXTE / "CorpusOne_output" / "corpus.md"
    if output.exists():
        output.unlink()
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_MIXTE),
            "-o",
            str(output),
            "--yes",
        ]
    )
    if code != 0:
        r.fail(f"Code {code}: {err}")
        return r
    if not output.exists():
        r.fail("Corpus non genere")
        return r
    content = output.read_text(encoding="utf-8")
    if "SOURCE:" not in content:
        r.fail("En-tete SOURCE manquant")
        return r
    r.ok("Corpus MD genere avec en-tetes SOURCE")
    return r


def test_dossier_mixte_pdf() -> Result:
    """S19.2 -- Meme dossier -> un PDF unique, texte selectionnable."""
    r = Result("Dossier mixte -> PDF")
    output = DOSSIER_MIXTE / "CorpusOne_output" / "corpus.pdf"
    if output.exists():
        output.unlink()
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_MIXTE),
            "-o",
            str(output),
            "--format",
            "pdf",
            "--yes",
        ]
    )
    if code != 0:
        r.fail(f"Code {code}: {err}")
        return r
    if not output.exists():
        r.fail("PDF non genere")
        return r
    if output.stat().st_size < 100:
        r.fail("PDF trop petit (vide ?)")
        return r
    r.ok("PDF genere")
    return r


def test_blocage_context() -> Result:
    """S19.2 -- TXT dont le compteur +15% > plafond -> pas de corpus, code 2."""
    r = Result("Blocage par plafond")
    output = DOSSIER_BLOCAGE / "corpus.md"
    if output.exists():
        output.unlink()
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_BLOCAGE),
            "-o",
            str(output),
            "--context",
            "100",
            "--yes",
        ]
    )
    if code != 2:
        r.fail(f"Code attendu: 2, obtenu: {code}")
        return r
    if output.exists():
        r.fail("Corpus ne devrait pas etre genere")
        return r
    r.ok("Blocage correct (code 2)")
    return r


def test_fichier_exe_ignore() -> Result:
    """S19.2 -- Fichier .exe dans le dossier -> ignore, present au rapport."""
    r = Result("Fichier .exe ignore")
    output = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_test.md"
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_MIXTE),
            "-o",
            str(output),
            "--dry-run",
            "--yes",
        ]
    )
    report = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_test_rapport.json"
    if not report.exists():
        r.fail("Rapport non genere")
        return r
    import json

    data = json.loads(report.read_text(encoding="utf-8"))
    ignored_names = [i["path"].split("/")[-1].split("\\")[-1] for i in data.get("ignored", [])]
    if "app.exe" not in ignored_names:
        r.fail(f"app.exe non liste dans ignores: {ignored_names}")
        return r
    r.ok("app.exe ignore et present au rapport")
    return r


def test_lock_file_ignore() -> Result:
    """S19.2 -- ~$w.docx -> ignore."""
    r = Result("Fichier verrou ~$ ignore")
    output = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_lock.md"
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_MIXTE),
            "-o",
            str(output),
            "--dry-run",
            "--yes",
        ]
    )
    report = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_lock_rapport.json"
    if not report.exists():
        r.fail("Rapport non genere")
        return r
    import json

    data = json.loads(report.read_text(encoding="utf-8"))
    ignored_names = [i["path"].split("/")[-1].split("\\")[-1] for i in data.get("ignored", [])]
    if "~$locked.docx" not in ignored_names:
        r.fail(f"~$locked.docx non liste: {ignored_names}")
        return r
    r.ok("~$locked.docx ignore")
    return r


def test_dry_run() -> Result:
    """S19.2 -- CLI --dry-run -> pas de corpus, rapport stats."""
    r = Result("Dry-run avec rapport")
    output = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_dry.md"
    code, out, err = run_cli(
        [
            "-i",
            str(DOSSIER_MIXTE),
            "-o",
            str(output),
            "--dry-run",
        ]
    )
    if output.exists():
        r.fail("Corpus ne devrait pas exister en dry-run")
        return r
    report_md = DOSSIER_MIXTE / "CorpusOne_output" / "corpus_dry_rapport.md"
    if not report_md.exists():
        r.fail("Rapport MD non genere en dry-run")
        return r
    r.ok("Dry-run: pas de corpus, rapport genere")
    return r


def test_list_formats() -> Result:
    """S6.3 -- --list-formats affiche les extensions."""
    r = Result("--list-formats")
    code, out, err = run_cli(["--list-formats"])
    if code != 0:
        r.fail(f"Code {code}")
        return r
    if ".pdf" not in out or ".docx" not in out:
        r.fail("Extensions manquantes")
        return r
    r.ok(f"Formats listes ({out.count('.')}) extensions)")
    return r


# ------------------------------------------------------------------------------
# Main
# ------------------------------------------------------------------------------


def main() -> int:
    _generate_test_files()
    print("=" * 60)
    print("Script de recette DocFuse / CorpusOne (CdC S21.4)")
    print("=" * 60)
    print()

    tests = [
        test_dossier_mixte_md,
        test_dossier_mixte_pdf,
        test_blocage_context,
        test_fichier_exe_ignore,
        test_lock_file_ignore,
        test_dry_run,
        test_list_formats,
    ]

    passed = 0
    failed = 0

    for test_fn in tests:
        result = test_fn()
        status = "PASS" if result.passed else "FAIL"
        print(f"  {status} -- {result.name}")
        if result.message:
            print(f"         {result.message}")
        if result.passed:
            passed += 1
        else:
            failed += 1
        print()

    print("=" * 60)
    print(f"Resultat: {passed} reussis, {failed} echoues sur {len(tests)} tests")
    print("=" * 60)

    return 0 if failed == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
