"""Génère les fichiers d'échantillon de test pour les formats binaires.

À exécuter une fois pour créer les fixtures :
    python tests/fixtures/generate_fixtures.py
"""

from __future__ import annotations

from email.message import EmailMessage
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def generate_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Titre Test Document", level=0)
    doc.add_paragraph(
        "Ceci est un paragraphe de test avec suffisamment de texte "
        "pour dépasser le seuil de quatre-vingt caractères et ne pas "
        "déclencher l alerte de pauvrete de texte."
    )
    doc.add_paragraph("Deuxieme paragraphe avec encore du texte.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Colonne A"
    table.cell(0, 1).text = "Colonne B"
    table.cell(1, 0).text = "Valeur 1"
    table.cell(1, 1).text = "Valeur 2"
    doc.save(str(FIXTURES_DIR / "sample.docx"))


def generate_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = "Titre Diapo"
    content_shape = slide.placeholders[1]
    content_shape.text = (
        "Contenu de test avec suffisamment de texte pour depasser "
        "le seuil de caracteres et ne pas declencher une alerte."
    )
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide2.shapes.title.text = "Diapo vide"
    prs.save(str(FIXTURES_DIR / "sample.pptx"))


def generate_xlsx() -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Feuille1"
    ws["A1"] = "Nom"
    ws["B1"] = "Valeur"
    ws["A2"] = "Test avec du texte suffisamment long pour eviter alerte"
    ws["B2"] = "42"
    ws["A3"] = "Ligne supplementaire"
    ws["B3"] = "99"
    ws2 = wb.create_sheet("Feuille2")
    ws2["A1"] = "Autre feuille avec du texte"
    ws2["A2"] = "Donnee supplementaire"
    wb.save(str(FIXTURES_DIR / "sample.xlsx"))


def generate_rtf() -> None:
    rtf_content = (
        r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}"
        r"\f0\fs24 Ceci est un document RTF de test avec suffisamment de texte "
        r"pour depasser le seuil de quatre-vingt caracteres et ne pas declencher "
        r"l alerte de pauvrete de texte.\par "
        r"Deuxieme paragraphe du document RTF.}"
    )
    (FIXTURES_DIR / "sample.rtf").write_text(rtf_content, encoding="ascii")


def generate_pdf() -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    doc = SimpleDocTemplate(
        str(FIXTURES_DIR / "sample.pdf"),
        pagesize=A4,
    )
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Titre du PDF de test", styles["Title"]),
        Spacer(1, 12),
        Paragraph(
            "Ceci est un paragraphe de test avec suffisamment de texte "
            "pour depasser le seuil de quatre-vingt caracteres et ne pas "
            "declencher l alerte de pauvrete de texte.",
            styles["Normal"],
        ),
        Spacer(1, 12),
        Paragraph("Deuxieme paragraphe du document PDF.", styles["Normal"]),
    ]
    doc.build(story)


def generate_eml() -> None:
    msg = EmailMessage()
    msg["Subject"] = "Test EML"
    msg["From"] = "sender@example.com"
    msg["To"] = "recipient@example.com"
    msg["Date"] = "Thu, 20 Aug 2026 12:00:00 +0000"
    msg.set_content(
        "Ceci est le corps de l email de test avec suffisamment de texte "
        "pour depasser le seuil de caracteres et ne pas declencher d alerte."
    )
    (FIXTURES_DIR / "sample.eml").write_bytes(bytes(msg))


def generate_odt() -> None:
    import zipfile

    # ODT minimal : content.xml + META-INF/manifest.xml + mimetype
    content_xml = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"
  xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"
  xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0"
  office:version="1.2">
  <office:body>
    <office:text>
      <text:p>Ceci est un paragraphe ODT de test avec suffisamment de texte pour eviter l alerte.</text:p>
      <text:p>Deuxieme paragraphe du document OpenDocument.</text:p>
    </office:text>
  </office:body>
</office:document-content>"""

    manifest_xml = """<?xml version="1.0" encoding="UTF-8"?>
<manifest:manifest xmlns:manifest="urn:oasis:names:tc:opendocument:xmlns:manifest:1.0">
  <manifest:file-entry manifest:full-path="/" manifest:media-type="application/vnd.oasis.opendocument.text"/>
  <manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>
</manifest:manifest>"""

    mimetype = b"application/vnd.oasis.opendocument.text"

    with zipfile.ZipFile(str(FIXTURES_DIR / "sample.odt"), "w") as zf:
        zf.writestr("mimetype", mimetype)
        zf.writestr("META-INF/manifest.xml", manifest_xml)
        zf.writestr("content.xml", content_xml)


if __name__ == "__main__":
    generate_docx()
    generate_pptx()
    generate_xlsx()
    generate_rtf()
    generate_pdf()
    generate_eml()
    generate_odt()
    print("All fixtures generated in", FIXTURES_DIR)
