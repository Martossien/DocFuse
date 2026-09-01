"""Tests de l'extracteur PPTX.

CdC §8.3 — Texte des shapes, tableaux, notes d'orateur.
D-107 — Graphiques, SmartArt, masque et dispositions ; marqueur de diapo qui
ne peut plus affirmer faussement qu'il n'y a rien à lire.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest

from docfuse.core.ocr.tesseract import TesseractEngine
from docfuse.extractors.pptx import PptxExtractor
from docfuse.models.file_status import FileStatus

_OCR_AVAILABLE = TesseractEngine().is_available()

# --- Fragments XML pour les cas que python-pptx ne sait pas écrire (D-107) ---
# SmartArt, forme posée sur un masque ou une disposition, contenu enveloppé
# dans <mc:AlternateContent> : rien de tout cela n'est exposé en écriture par
# la bibliothèque, il faut descendre dans le XML du paquet.

_NS_DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"

_CT_DIAGRAM = (
    '<Override PartName="/ppt/diagrams/data1.xml" ContentType="application/'
    'vnd.openxmlformats-officedocument.drawingml.diagramData+xml"/>'
)

_REL_DIAGRAM = (
    f'<Relationship Id="rIdDgm1" Type="{_NS_R}/diagramData" Target="../diagrams/data1.xml"/>'
)

_GRAPHIC_FRAME_DIAGRAM = (
    "<p:graphicFrame><p:nvGraphicFramePr>"
    '<p:cNvPr id="42" name="Diagramme 1"/><p:cNvGraphicFramePr/><p:nvPr/>'
    "</p:nvGraphicFramePr>"
    '<p:xfrm><a:off x="838200" y="1825625"/><a:ext cx="5486400" cy="3200400"/></p:xfrm>'
    f'<a:graphic><a:graphicData uri="{_NS_DGM}">'
    f'<dgm:relIds xmlns:dgm="{_NS_DGM}" xmlns:r="{_NS_R}" '
    'r:dm="rIdDgm1" r:lo="rIdDgm1" r:qs="rIdDgm1" r:cs="rIdDgm1"/>'
    "</a:graphicData></a:graphic></p:graphicFrame>"
)

_GRAPHIC_FRAME_INCONNU = (
    "<p:graphicFrame><p:nvGraphicFramePr>"
    '<p:cNvPr id="43" name="Objet mystere"/><p:cNvGraphicFramePr/><p:nvPr/>'
    "</p:nvGraphicFramePr>"
    '<p:xfrm><a:off x="838200" y="1825625"/><a:ext cx="5486400" cy="3200400"/></p:xfrm>'
    '<a:graphic><a:graphicData uri="http://exemple.invalid/format-inconnu"/></a:graphic>'
    "</p:graphicFrame>"
)


def _data_model(*textes: str) -> str:
    """Modèle de données SmartArt (`ppt/diagrams/data1.xml`) minimal."""
    points = "".join(
        f'<dgm:pt modelId="{{p{i}}}"><dgm:prSet/><dgm:spPr/><dgm:t>'
        "<a:bodyPr/><a:lstStyle/><a:p><a:r>"
        f'<a:rPr lang="fr-FR"/><a:t>{texte}</a:t></a:r></a:p></dgm:t></dgm:pt>'
        for i, texte in enumerate(textes)
    )
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<dgm:dataModel xmlns:dgm="{_NS_DGM}" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<dgm:ptLst><dgm:pt modelId="{{doc}}" type="doc"><dgm:prSet/></dgm:pt>{points}</dgm:ptLst>'
        "<dgm:cxnLst/><dgm:bg/><dgm:whole/></dgm:dataModel>"
    )


def _textbox_xml(shape_id: int, nom: str, texte: str) -> str:
    """Zone de texte `<p:sp>` autonome (pas un espace réservé)."""
    return (
        "<p:sp><p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="{nom}"/><p:cNvSpPr txBox="1"/><p:nvPr/>'
        "</p:nvSpPr>"
        '<p:spPr><a:xfrm><a:off x="457200" y="6172200"/>'
        '<a:ext cx="5486400" cy="365125"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        "<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r>"
        f'<a:rPr lang="fr-FR"/><a:t>{texte}</a:t>'
        "</a:r></a:p></p:txBody></p:sp>"
    )


def _alternate_content(texte: str) -> str:
    """`<mc:AlternateContent>` : les deux branches portent le même texte —
    l'extracteur ne doit en lire qu'une."""
    return (
        f'<mc:AlternateContent xmlns:mc="{_NS_MC}">'
        '<mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"'
        ' Requires="p14">' + _textbox_xml(50, "Choix", texte) + "</mc:Choice>"
        "<mc:Fallback>" + _textbox_xml(51, "Repli", texte) + "</mc:Fallback>"
        "</mc:AlternateContent>"
    )


def _reecrire_zip(
    chemin: Path,
    remplacements: dict[str, tuple[str, str]],
    ajouts: dict[str, str] | None = None,
) -> None:
    """Réécrit un .pptx en place : substitutions dans des parties existantes
    et parties ajoutées."""
    source = chemin.with_suffix(".source.pptx")
    chemin.replace(source)
    with (
        zipfile.ZipFile(source) as src,
        zipfile.ZipFile(chemin, "w", zipfile.ZIP_DEFLATED) as dst,
    ):
        for info in src.infolist():
            donnees = src.read(info.filename)
            if info.filename in remplacements:
                cible, remplacant = remplacements[info.filename]
                texte = donnees.decode("utf-8")
                assert cible in texte, f"{info.filename}: {cible!r} introuvable"
                donnees = texte.replace(cible, remplacant).encode("utf-8")
            dst.writestr(info, donnees)
        for nom, contenu in (ajouts or {}).items():
            dst.writestr(nom, contenu)
    source.unlink()


def _deck_vierge(tmp_path: Path, name: str, nb_diapos: int = 1) -> Path:
    """Présentation avec `nb_diapos` diapositives vides (disposition Blank)."""
    from pptx import Presentation

    f = tmp_path / name
    prs = Presentation()
    for _ in range(nb_diapos):
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(f))
    return f


def _pptx_with_image(tmp_path: Path, name: str, image_bytes: bytes) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    f = tmp_path / name
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(image_bytes), Inches(1), Inches(1), Inches(2), Inches(2))
    prs.save(str(f))
    return f


def _red_square_png() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (50, 50), "red").save(buf, format="PNG")
    return buf.getvalue()


class TestPptxExtractor:
    def test_extract_basic_slide(self, tmp_path: Path) -> None:
        from pptx import Presentation

        f = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Titre Diapo"
        slide.placeholders[1].text = (
            "Contenu de la diapo avec suffisamment de texte "
            "pour depasser le seuil de quatre-vingt caracteres."
        )
        prs.save(str(f))

        result = PptxExtractor.extract(f, "test.pptx")
        assert result.status is FileStatus.READY
        assert "Titre Diapo" in result.text
        assert result.page_count == 1

    def test_empty_slide_marker(self, tmp_path: Path) -> None:
        """Une diapo réellement dépourvue de forme porteuse de texte reste
        signalée par un marqueur `[[DIAPO 1: ...]]`.

        D-107 — le libellé attendu a changé, volontairement. L'ancien
        `[[DIAPO 1: aucun texte extractible]]` était une affirmation sur le
        *contenu de la diapo* alors que l'extracteur n'avait regardé que
        `has_text_frame` et `has_table` : un graphique, un SmartArt ou une
        image sans OCR déclenchaient le même marqueur, qui envoyait
        l'auditeur conclure « diapo image, rien à lire » sur un organigramme
        nominatif. Le nouveau libellé décrit ce qui a été inspecté, pas ce
        que la diapo contiendrait. Le test verrouille désormais la forme du
        marqueur et surtout l'absence de l'ancienne affirmation."""
        f = tmp_path / "empty.pptx"
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
        prs.save(str(f))

        result = PptxExtractor.extract(f, "empty.pptx")
        assert "[[DIAPO 1:" in result.text
        assert "aucune forme porteuse de texte" in result.text
        assert "aucun texte extractible" not in result.text

    def test_accepts(self) -> None:
        assert PptxExtractor.accepts(Path("test.pptx")) is True
        assert PptxExtractor.accepts(Path("test.docx")) is False

    def test_safe_extract_no_crash(self, tmp_path: Path) -> None:
        f = tmp_path / "nonexistent.pptx"
        result = PptxExtractor.safe_extract(f, "nonexistent.pptx")
        assert result.status is FileStatus.ERROR

    def test_password_protected_gives_clear_error(self, tmp_path: Path) -> None:
        """D-089 : un .pptx protégé par mot de passe à l'ouverture (conteneur
        OLE2, plus un ZIP) donnait un `PackageNotFoundError` bas niveau,
        jamais un message disant à l'utilisateur que le fichier est protégé."""
        f = tmp_path / "protected.pptx"
        f.write_bytes(bytes((0xD0, 0xCF, 0x11, 0xE0, 0xA1, 0xB1, 0x1A, 0xE1)) + b"\x00" * 500)

        result = PptxExtractor.extract(f, "protected.pptx")
        assert result.status is FileStatus.ERROR
        assert result.error_message is not None
        assert "mot de passe" in result.error_message.lower()

    def test_fixture_file(self) -> None:
        fixture = Path(__file__).resolve().parent.parent / "fixtures" / "sample.pptx"
        if fixture.exists():
            result = PptxExtractor.extract(fixture, "sample.pptx")
            assert result.status is FileStatus.READY
            assert "Diapo" in result.text

    def test_grouped_shapes_text_is_extracted(self, tmp_path: Path) -> None:
        """D-074 : le texte dans une forme groupée (GroupShape — schémas,
        diagrammes annotés) ne doit pas disparaître. shape.has_text_frame
        renvoie False pour le conteneur groupe lui-même ; sans récursion
        dans shape.shapes, tout son contenu est invisible."""
        from pptx import Presentation
        from pptx.util import Inches

        f = tmp_path / "grouped.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tb1.text_frame.text = "TEXTE_DANS_GROUPE_1"
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
        tb2.text_frame.text = "TEXTE_DANS_GROUPE_2"
        slide.shapes.add_group_shape([tb1, tb2])
        prs.save(str(f))

        result = PptxExtractor.extract(f, "grouped.pptx")
        assert result.status is FileStatus.READY
        assert "TEXTE_DANS_GROUPE_1" in result.text
        assert "TEXTE_DANS_GROUPE_2" in result.text

    def test_embedded_image_export_creates_embedded_images(self, tmp_path: Path) -> None:
        """D-091 : export actif -> l'image est capturée avec un nom explicite
        (incluant le numéro de diapo) et un tag `[[IMAGE: ...]]` est inséré
        au point d'apparition."""
        f = _pptx_with_image(tmp_path, "with_image.pptx", _red_square_png())

        result = PptxExtractor.extract(f, "with_image.pptx", extract_images=True)
        assert result.status is FileStatus.READY
        assert len(result.embedded_images) == 1
        image = result.embedded_images[0]
        assert image.filename.startswith("with_image__slide1__img1")
        assert image.data
        assert f"[[IMAGE: {image.filename}]]" in result.text

    def test_embedded_image_export_disabled_by_default(self, tmp_path: Path) -> None:
        """D-091 : sans `extract_images`, aucune image capturée, texte inchangé
        (non-régression stricte vis-à-vis du comportement pré-D-091)."""
        f = _pptx_with_image(tmp_path, "with_image2.pptx", _red_square_png())

        result = PptxExtractor.extract(f, "with_image2.pptx")
        assert result.embedded_images == []
        assert "[[IMAGE" not in result.text

    @pytest.mark.skipif(not _OCR_AVAILABLE, reason="Tesseract non installé")
    def test_embedded_image_ocr_extracts_text_automatically(self, tmp_path: Path) -> None:
        """D-091 : corrige le bug signalé par l'utilisateur (PPTX où le texte
        est dans une image) — OCR automatique, sans avoir besoin d'activer
        l'export."""
        from PIL import Image, ImageDraw

        img = Image.new("RGB", (600, 150), "white")
        draw = ImageDraw.Draw(img)
        draw.text((10, 50), "Bonjour le monde", fill="black")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        f = _pptx_with_image(tmp_path, "with_text_image.pptx", buf.getvalue())

        result = PptxExtractor.extract(f, "with_text_image.pptx")
        assert "onjour" in result.text or "monde" in result.text
        assert result.embedded_images == []


class TestPptxContenuPerduD107:
    """D-107 — contenu qu'un `<p:graphicFrame>`, un masque ou une disposition
    portaient et que l'extracteur jetait, en écrivant par-dessus
    « aucun texte extractible »."""

    def test_texte_du_graphique_est_extrait(self, tmp_path: Path) -> None:
        """Un graphique est un `<p:graphicFrame>` : ni `has_text_frame` ni
        `has_table`. Titre, nom de série et catégories — « Salaires »,
        « Budget RH 2024 » — sortaient du corpus en silence."""
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        f = tmp_path / "graphique.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        donnees = CategoryChartData()
        donnees.categories = ["CAT_Salaires", "CAT_Primes"]
        donnees.add_series("SERIE_Budget_RH_2024", (52000.0, 8000.0))
        cadre = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(8), Inches(5), donnees
        )
        cadre.chart.has_title = True
        cadre.chart.chart_title.text_frame.text = "TITRE_Repartition_Masse_Salariale"
        prs.save(str(f))

        result = PptxExtractor.extract(f, "graphique.pptx")
        assert result.status is FileStatus.READY
        assert "TITRE_Repartition_Masse_Salariale" in result.text
        assert "SERIE_Budget_RH_2024" in result.text
        assert "CAT_Salaires" in result.text
        assert "CAT_Primes" in result.text
        # La diapo n'est plus déclarée muette.
        assert "[[DIAPO 1:" not in result.text
        # Les références de cellules du classeur lié ne sont pas du contenu.
        assert "Sheet1!" not in result.text

    def test_texte_du_smartart_est_extrait(self, tmp_path: Path) -> None:
        """Un SmartArt vit dans `ppt/diagrams/dataN.xml`, atteint par le
        `<dgm:relIds r:dm>` du graphicFrame. python-pptx n'en expose rien —
        un organigramme est pourtant du nominatif pur."""
        f = _deck_vierge(tmp_path, "smartart.pptx")
        _reecrire_zip(
            f,
            {
                "ppt/slides/slide1.xml": ("</p:spTree>", _GRAPHIC_FRAME_DIAGRAM + "</p:spTree>"),
                "ppt/slides/_rels/slide1.xml.rels": (
                    "</Relationships>",
                    _REL_DIAGRAM + "</Relationships>",
                ),
                "[Content_Types].xml": ("</Types>", _CT_DIAGRAM + "</Types>"),
            },
            {
                "ppt/diagrams/data1.xml": _data_model(
                    "SMARTART_Organigramme_Direction_RH", "SMARTART_Sophie_Martin_DRH"
                )
            },
        )

        result = PptxExtractor.extract(f, "smartart.pptx")
        assert result.status is FileStatus.READY
        assert "SMARTART_Organigramme_Direction_RH" in result.text
        assert "SMARTART_Sophie_Martin_DRH" in result.text
        assert "[[DIAPO 1:" not in result.text

    def test_smartart_orphelin_est_recupere(self, tmp_path: Path) -> None:
        """Un SmartArt posé sur un masque, une disposition ou une page de
        notes n'est réclamé par aucune diapo. Le balayage de fin de parcours
        le restitue dans sa propre section plutôt que de le perdre."""
        f = _deck_vierge(tmp_path, "orphelin.pptx")
        _reecrire_zip(
            f,
            {"[Content_Types].xml": ("</Types>", _CT_DIAGRAM + "</Types>")},
            {"ppt/diagrams/data1.xml": _data_model("SMARTART_ORPHELIN_Comite_Direction")},
        )

        result = PptxExtractor.extract(f, "orphelin.pptx")
        assert "SMARTART_ORPHELIN_Comite_Direction" in result.text
        assert "## Diagrammes hors diapositive" in result.text

    def test_smartart_rattache_nest_pas_repete_en_fin_de_corpus(self, tmp_path: Path) -> None:
        """Le balayage des diagrammes orphelins ne doit pas redonner ce qu'une
        diapo a déjà rendu — sinon chaque SmartArt compterait deux fois."""
        f = _deck_vierge(tmp_path, "smartart_unique.pptx")
        _reecrire_zip(
            f,
            {
                "ppt/slides/slide1.xml": ("</p:spTree>", _GRAPHIC_FRAME_DIAGRAM + "</p:spTree>"),
                "ppt/slides/_rels/slide1.xml.rels": (
                    "</Relationships>",
                    _REL_DIAGRAM + "</Relationships>",
                ),
                "[Content_Types].xml": ("</Types>", _CT_DIAGRAM + "</Types>"),
            },
            {"ppt/diagrams/data1.xml": _data_model("SMARTART_Sophie_Martin_DRH")},
        )

        result = PptxExtractor.extract(f, "smartart_unique.pptx")
        assert result.text.count("SMARTART_Sophie_Martin_DRH") == 1
        assert "## Diagrammes hors diapositive" not in result.text

    def test_texte_du_masque_est_extrait(self, tmp_path: Path) -> None:
        """Un bandeau de classification est posé sur le masque précisément
        pour valoir sur toutes les diapositives. `slide.shapes` ne le voit
        pas."""
        f = _deck_vierge(tmp_path, "masque.pptx")
        _reecrire_zip(
            f,
            {
                "ppt/slideMasters/slideMaster1.xml": (
                    "</p:spTree>",
                    _textbox_xml(90, "Bandeau", "MASQUE_bandeau_DIFFUSION_RESTREINTE")
                    + "</p:spTree>",
                )
            },
        )

        result = PptxExtractor.extract(f, "masque.pptx")
        assert "MASQUE_bandeau_DIFFUSION_RESTREINTE" in result.text
        assert "## Gabarit" in result.text

    def test_texte_de_la_disposition_est_extrait(self, tmp_path: Path) -> None:
        """Même perte pour une mention posée sur la disposition (slideLayout)."""
        f = _deck_vierge(tmp_path, "disposition.pptx")
        # slide_layouts[6] « Blank » == ppt/slideLayouts/slideLayout7.xml
        _reecrire_zip(
            f,
            {
                "ppt/slideLayouts/slideLayout7.xml": (
                    "</p:spTree>",
                    _textbox_xml(91, "Mention", "DISPOSITION_mention_Projet_Confidentiel")
                    + "</p:spTree>",
                )
            },
        )

        result = PptxExtractor.extract(f, "disposition.pptx")
        assert "DISPOSITION_mention_Projet_Confidentiel" in result.text
        assert "diapos 1" in result.text or "toutes les diapositives" in result.text

    def test_texte_du_masque_nest_ecrit_quune_fois(self, tmp_path: Path) -> None:
        """Le texte du masque est commun à toutes les diapos : le répéter sur
        chacune gonflerait le corpus et fausserait le comptage de jetons. Il
        est écrit une seule fois, en section « Gabarit »."""
        f = _deck_vierge(tmp_path, "masque_x5.pptx", nb_diapos=5)
        _reecrire_zip(
            f,
            {
                "ppt/slideMasters/slideMaster1.xml": (
                    "</p:spTree>",
                    _textbox_xml(90, "Bandeau", "MASQUE_bandeau_DIFFUSION_RESTREINTE")
                    + "</p:spTree>",
                )
            },
        )

        result = PptxExtractor.extract(f, "masque_x5.pptx")
        assert result.page_count == 5
        assert result.text.count("MASQUE_bandeau_DIFFUSION_RESTREINTE") == 1
        # ...et pas à l'intérieur d'une diapo : la section « Gabarit » précède
        # la première `## Diapo`.
        assert result.text.index("MASQUE_bandeau_DIFFUSION_RESTREINTE") < result.text.index(
            "## Diapo 1"
        )

    def test_marques_de_gabarit_absentes_du_corpus(self, tmp_path: Path) -> None:
        """Le masque du modèle par défaut porte le texte d'invite des espaces
        réservés (« Click to edit Master title style »), la date figée du
        modèle et le littéral « ‹#› » du numéro de diapo. Ce n'est pas du
        contenu : lire le masque ne doit pas remplir le corpus de bruit."""
        f = _deck_vierge(tmp_path, "bruit.pptx")
        _reecrire_zip(
            f,
            {
                "ppt/slideMasters/slideMaster1.xml": (
                    "</p:spTree>",
                    _textbox_xml(90, "Bandeau", "MASQUE_utile") + "</p:spTree>",
                )
            },
        )

        result = PptxExtractor.extract(f, "bruit.pptx")
        assert "MASQUE_utile" in result.text
        assert "Click to edit" not in result.text
        assert "Master text styles" not in result.text
        assert "Second level" not in result.text
        assert "1/27/13" not in result.text
        assert "‹#›" not in result.text

    def test_contenu_alternatif_est_extrait_une_seule_fois(self, tmp_path: Path) -> None:
        """python-pptx saute sans un mot les éléments de `<p:spTree>` qu'il ne
        connaît pas — au premier rang `<mc:AlternateContent>`, l'enveloppe de
        l'encre manuscrite et des nouveautés PowerPoint. On lit une seule des
        deux branches : `mc:Choice` et `mc:Fallback` portent le même texte."""
        f = _deck_vierge(tmp_path, "altcontent.pptx")
        _reecrire_zip(
            f,
            {
                "ppt/slides/slide1.xml": (
                    "</p:spTree>",
                    _alternate_content("ALTCONTENT_note_Sophie_Martin") + "</p:spTree>",
                )
            },
        )

        result = PptxExtractor.extract(f, "altcontent.pptx")
        assert result.text.count("ALTCONTENT_note_Sophie_Martin") == 1


class TestPptxMarqueurHonnete:
    """D-107 — le marqueur de diapo ne doit jamais affirmer faussement qu'il
    n'y a rien à lire."""

    def test_objet_graphique_inconnu_est_nomme_pas_tu(self, tmp_path: Path) -> None:
        """Un `<p:graphicFrame>` d'un type que DocFuse ne sait pas lire :
        avant, la diapo sortait avec « aucun texte extractible », ce qui
        orientait l'auditeur vers « diapo image, rien à lire ». Le marqueur
        doit dire que quelque chose n'a pas été analysé."""
        f = _deck_vierge(tmp_path, "inconnu.pptx")
        _reecrire_zip(
            f,
            {"ppt/slides/slide1.xml": ("</p:spTree>", _GRAPHIC_FRAME_INCONNU + "</p:spTree>")},
        )

        result = PptxExtractor.extract(f, "inconnu.pptx")
        assert "aucun texte extractible" not in result.text
        assert "non analysé par DocFuse" in result.text
        assert "format-inconnu" in result.text

    def test_image_sans_ocr_est_signalee(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """Sans moteur OCR et sans export, une image sortait du corpus sans
        laisser de trace et la diapo était déclarée « aucun texte
        extractible » — le mensonge le plus fréquent, `DocFuse.exe` standard
        n'embarquant pas Tesseract. Cas d'une diapo qui ne porte QUE l'image :
        le marqueur doit être complet."""
        monkeypatch.setattr("docfuse.extractors.pptx.resolve_ocr_engine", lambda: None)
        f = _pptx_with_image(tmp_path, "photo.pptx", _red_square_png())

        result = PptxExtractor.extract(f, "photo.pptx")
        assert "aucun texte extractible" not in result.text
        assert "image non analysée" in result.text

    def test_image_sans_ocr_nest_pas_signalee_si_la_diapo_porte_du_texte(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """Arbitrage de volume, verrouillé exprès (D-107).

        Une diapo qui porte déjà son texte n'a pas besoin d'une ligne de plus
        pour dire qu'elle contient aussi une photo : le fichier est déjà
        marqué `FileStatus.IMAGES` par `image_count`. Mesuré sur huit decks
        réels sans Tesseract, la mention par diapo ajoutait ~8 % de corpus —
        le gonflement même que l'on refuse au texte du masque. Ce test
        empêche de la réintroduire par mégarde."""
        from pptx import Presentation
        from pptx.util import Inches

        monkeypatch.setattr("docfuse.extractors.pptx.resolve_ocr_engine", lambda: None)
        f = _pptx_with_image(tmp_path, "photo_et_texte.pptx", _red_square_png())
        prs = Presentation(str(f))
        zone = prs.slides[0].shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
        zone.text_frame.text = "TITRE_Resultats_2024"
        prs.save(str(f))

        result = PptxExtractor.extract(f, "photo_et_texte.pptx")
        assert "TITRE_Resultats_2024" in result.text
        assert "[[DIAPO 1:" not in result.text

    def test_diapo_avec_texte_signale_aussi_ce_qui_na_pas_ete_lu(self, tmp_path: Path) -> None:
        """Le signalement ne doit pas dépendre du fait que la diapo soit vide :
        un titre anodin à côté d'un objet illisible ne doit pas laisser croire
        que le titre est tout ce qu'il y a."""
        from pptx import Presentation
        from pptx.util import Inches

        f = tmp_path / "mixte.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        zone = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        zone.text_frame.text = "TITRE_Annexe"
        prs.save(str(f))
        _reecrire_zip(
            f,
            {"ppt/slides/slide1.xml": ("</p:spTree>", _GRAPHIC_FRAME_INCONNU + "</p:spTree>")},
        )

        result = PptxExtractor.extract(f, "mixte.pptx")
        assert "TITRE_Annexe" in result.text
        assert "non analysé par DocFuse" in result.text


class TestPptxCdC83:
    """CdC §8.3 — tableaux et notes d'orateur. Ces deux extractions
    n'étaient couvertes par aucun test : les supprimer entièrement de
    `pptx.py` ne faisait tomber aucune assertion (trou de mutation relevé
    par l'audit D-107), alors que les notes d'orateur sont en gras au cahier
    des charges."""

    def test_tableau_de_diapo_est_extrait(self, tmp_path: Path) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        f = tmp_path / "tableau.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
        table.cell(0, 0).text = "TAB_Nom"
        table.cell(0, 1).text = "TAB_Salaire"
        table.cell(1, 0).text = "TAB_Sophie_Martin"
        table.cell(1, 1).text = "TAB_52000_EUR"
        prs.save(str(f))

        result = PptxExtractor.extract(f, "tableau.pptx")
        assert result.status is FileStatus.READY
        assert "TAB_Nom | TAB_Salaire" in result.text
        assert "TAB_Sophie_Martin | TAB_52000_EUR" in result.text

    def test_notes_dorateur_sont_extraites(self, tmp_path: Path) -> None:
        from pptx import Presentation

        f = tmp_path / "notes.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.notes_slide.notes_text_frame.text = "NOTES_licenciement_economique_Q3"
        prs.save(str(f))

        result = PptxExtractor.extract(f, "notes.pptx")
        assert result.status is FileStatus.READY
        assert "[Notes] NOTES_licenciement_economique_Q3" in result.text
